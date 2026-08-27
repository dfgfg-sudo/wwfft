import os, sys
from pathlib import Path
from datetime import datetime

# ----- 首次运行需安装依赖（命令行执行）-----
# pip install python-docx PyPDF2 openpyxl python-pptx

import docx, PyPDF2, openpyxl
from pptx import Presentation

# ========== 只改这里 ==========
ROOT_DIR = r"D:/你的总文件夹"      # 您存放多个文件夹的父目录
OUTPUT_FILE = r"./完整知识库.md"    # 输出的知识库文件
# ==============================

def extract_text(file_path):
    """万能提取器：任何文件都能吐出文本或占位，决不出错"""
    ext = Path(file_path).suffix.lower()
    try:
        # 纯文本 / 代码（强制UTF-8，坏字符替换为�）
        if ext in ['.txt','.md','.csv','.json','.xml','.html','.log',
                   '.py','.js','.java','.c','.cpp','.sh','.sql']:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
        # Word
        elif ext == '.docx':
            return '\n'.join([p.text for p in docx.Document(file_path).paragraphs])
        # PDF
        elif ext == '.pdf':
            reader = PyPDF2.PdfReader(file_path)
            text = ''.join([p.extract_text() or '' for p in reader.pages])
            return text if text.strip() else "[扫描件PDF，无原生文字]"
        # Excel
        elif ext in ['.xlsx','.xls']:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            all_txt = ''
            for sheet in wb.worksheets:
                all_txt += f"\n--- 工作表 {sheet.title} ---\n"
                for row in sheet.iter_rows(values_only=True):
                    row_str = ' '.join(str(c) for c in row if c is not None)
                    if row_str: all_txt += row_str + '\n'
            return all_txt
        # PPT
        elif ext == '.pptx':
            prs = Presentation(file_path)
            return '\n'.join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape,'text')])
        else:
            return "[占位：二进制/特殊格式文件]"
    except Exception as e:
        return f"[提取失败：{str(e)}]"

if __name__ == "__main__":
    if not os.path.exists(ROOT_DIR):
        print(f"错误：路径 {ROOT_DIR} 不存在")
        sys.exit(1)
    
    total = success = 0
    with open(OUTPUT_FILE, 'w', encoding='utf-8') as out:
        out.write(f"# 📚 完整知识库（全目录镜像）\n生成时间：{datetime.now()}\n根目录：`{ROOT_DIR}`\n\n---\n\n")
        # 遍历（强制排序，保证顺序）
        for dirpath, _, files in sorted(os.walk(ROOT_DIR)):
            rel = os.path.relpath(dirpath, ROOT_DIR)
            level = 1 if rel == '.' else rel.count(os.sep) + 2
            out.write(f"{'#'*level} 📁 {rel if rel != '.' else '根目录'}\n\n")
            for file in sorted(files):
                total += 1
                full = os.path.join(dirpath, file)
                out.write(f"### 📄 {file}\n路径：`{full}`\n\n")
                content = extract_text(full)
                if not content.startswith('[') and not content.startswith('提取失败'):
                    success += 1
                out.write(content + "\n\n---\n\n")
        # 统计校验
        out.write(f"\n## ✅ 完整性校验\n- 总文件数：{total}\n- 成功提取文本：{success}\n- 占位/修复标记：{total-success}\n")
    print(f"知识库已生成：{OUTPUT_FILE}\n总文件数：{total}，成功提取：{success}")